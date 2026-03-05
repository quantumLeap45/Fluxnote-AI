import csv
import io
from typing import Optional

ALLOWED_EXTENSIONS = {"pdf", "docx", "txt", "csv", "pptx"}

_CONTENT_TYPE_MAP = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "text/plain": "txt",
    "text/csv": "csv",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}


def _resolve_format(content_type: str, filename: str) -> Optional[str]:
    """Determine file format from content_type, then fall back to file extension."""
    fmt = _CONTENT_TYPE_MAP.get(content_type.strip().lower())
    if fmt:
        return fmt
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    return ext if ext in ALLOWED_EXTENSIONS else None


def _parse_pdf(content: bytes) -> Optional[str]:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=content, filetype="pdf")
        pages = [page.get_text() for page in doc]
        return "\n".join(pages)
    except Exception:
        return None


def _parse_docx(content: bytes) -> Optional[str]:
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        parts = []

        # Paragraphs (body text)
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)

        # Tables — rubrics, marking criteria, structured content
        for table in doc.tables:
            for row in table.rows:
                seen = set()
                unique_cells = []
                for cell in row.cells:
                    if cell._tc not in seen:
                        seen.add(cell._tc)
                        unique_cells.append(cell)
                row_text = " | ".join(
                    cell.text.strip() for cell in unique_cells if cell.text.strip()
                )
                if row_text:
                    parts.append(row_text)

        # Headers and footers — often contain module code, due date, title
        for section in doc.sections:
            for container in (section.header, section.footer):
                if not container.is_linked_to_previous:
                    for para in container.paragraphs:
                        if para.text.strip():
                            parts.append(para.text)

        return "\n".join(parts) or None
    except Exception:
        return None


def _parse_txt(content: bytes) -> Optional[str]:
    try:
        return content.decode("utf-8", errors="replace")
    except Exception:
        return None


def _parse_csv(content: bytes) -> Optional[str]:
    try:
        text = content.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text))
        rows = [", ".join(row) for row in reader]
        return "\n".join(rows)
    except Exception:
        return None


def _parse_pptx(content: bytes) -> Optional[str]:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return None


_PARSERS = {
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "txt": _parse_txt,
    "csv": _parse_csv,
    "pptx": _parse_pptx,
}


def parse_file(content: bytes, content_type: str, filename: str) -> Optional[str]:
    """Parse uploaded file bytes into plain text.

    Returns extracted text on success, or None if the format is unsupported
    or parsing fails.
    """
    fmt = _resolve_format(content_type, filename)
    if fmt is None:
        return None
    parser = _PARSERS.get(fmt)
    if parser is None:
        return None
    return parser(content)
